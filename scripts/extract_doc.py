#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
鳍点AI 文档附件文本提取器
支持: .xlsx/.xlsm (openpyxl)  .docx (zip+xml)  .pptx (python-pptx)  .pdf (pypdf)
用法: python extract_doc.py <文件路径>
输出: 提取的纯文本到 stdout（UTF-8）
"""
import sys
import os
import io
import zipfile
import re


def extract_docx(path):
    """docx 本质是 zip，直接解析 word/document.xml 中的 <w:t>"""
    with zipfile.ZipFile(path) as z:
        xml = z.read("word/document.xml").decode("utf-8", "replace")
    # 段落分隔：</w:p> 换行
    xml = re.sub(r"</w:p>", "\n", xml)
    # 表格单元格分隔
    xml = re.sub(r"</w:tc>", " | ", xml)
    texts = re.findall(r"<w:t[^>]*>([^<]*)</w:t>", xml)
    out = "".join(texts)
    # 上面的 findall 会丢掉换行标记，改为按标签流拼接更稳：
    out = _docx_flow(xml)
    return out.strip()


def _docx_flow(xml):
    parts = []
    pos = 0
    token = re.compile(r"<w:t[^>]*>(.*?)</w:t>|<w:p[ />]|</w:p>|<w:tab[ />]|<w:br[ />]|</w:tr>|</w:tc>")
    for m in token.finditer(xml):
        if m.group(1) is not None:
            parts.append(m.group(1))
        elif m.group(0).startswith("<w:p") or m.group(0) == "</w:p>":
            parts.append("\n")
        elif m.group(0).startswith("</w:tc>"):
            parts.append(" | ")
        elif m.group(0).startswith("</w:tr>"):
            parts.append("\n")
        elif m.group(0).startswith("<w:tab") or m.group(0).startswith("<w:br"):
            parts.append(" ")
    return "".join(parts)


def extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append("【工作表】" + ws.title)
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(vals):
                lines.append(" | ".join(vals).rstrip(" |"))
    wb.close()
    return "\n".join(lines)


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append("【幻灯片 %d】" % i)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs)
                    if t.strip():
                        lines.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.replace("\n", " ") for c in row.cells]
                    lines.append(" | ".join(cells))
    return "\n".join(lines)


def extract_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path)
    lines = []
    for i, page in enumerate(reader.pages, 1):
        t = page.extract_text() or ""
        if t.strip():
            lines.append("【第 %d 页】" % i)
            lines.append(t.strip())
    return "\n".join(lines)


def main():
    if len(sys.argv) < 2:
        print("用法: python extract_doc.py <文件路径>", file=sys.stderr)
        return 2
    path = sys.argv[1]
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            out = extract_docx(path)
        elif ext in (".xlsx", ".xlsm"):
            out = extract_xlsx(path)
        elif ext == ".pptx":
            out = extract_pptx(path)
        elif ext == ".pdf":
            out = extract_pdf(path)
        else:
            print("不支持的类型: " + ext, file=sys.stderr)
            return 2
    except Exception as e:
        print("解析失败: %s: %s" % (type(e).__name__, e), file=sys.stderr)
        return 1
    sys.stdout.buffer.write(out.encode("utf-8", "replace"))
    return 0


if __name__ == "__main__":
    sys.exit(main())
